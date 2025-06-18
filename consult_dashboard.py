# consult_dashboard.py
import json
import os
import base64
import streamlit as st
import streamlit.components.v1 as components
import re
import textwrap
import gspread
from google.oauth2.service_account import Credentials
from datetime import datetime
from streamlit.components.v1 import html
import pandas as pd
from gspread_dataframe import set_with_dataframe

from openai import OpenAI

import io
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas


from dotenv import load_dotenv
load_dotenv()


from models.foot_traffic_model import run as run_foot_traffic
from models.margin_model import run as run_margin
from models.arcgis_explorer import *
from models.deep_dive_model import run_deep_dive
from models.optimizer import run_optimization
from models.presentation import generate_impress_html


# --- GOOGLE SHEETS SETUP ---
SHEET_ID    = "12Qvpi5jOdtWRaa1aL6yglCAJ5tFphW1fHsF8apTlEV4"
WS_NAME     = "Data"
AUTH_SCOPES = ["https://www.googleapis.com/auth/spreadsheets"]

# 1) Decode the base64 JSON
b64 = os.environ["GCLOUD_SA_KEY_B64"]
sa_json = base64.b64decode(b64).decode("utf-8")
info = json.loads(sa_json)

# 2) Build your creds & client
SCOPES = ["https://www.googleapis.com/auth/spreadsheets"]

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
client = OpenAI(api_key=OPENAI_API_KEY)


@st.cache_resource(show_spinner=False)
def get_gspread_client():
    creds  = Credentials.from_service_account_info(info, scopes=AUTH_SCOPES)
    return gspread.authorize(creds)

@st.cache_resource(show_spinner=False)
def get_worksheet():
    client = get_gspread_client()
    return client.open_by_key(SHEET_ID).worksheet(WS_NAME)

@st.cache_resource(show_spinner=False)
def get_pos_sheet():
    client = get_gspread_client()
    try:
        return client.open_by_key(SHEET_ID).worksheet("POS Raw")
    except gspread.WorksheetNotFound:
        return None

ws = get_worksheet()
# Cache the header row once
@st.cache_data(show_spinner=False)
def get_headers():
    return [h for h in ws.row_values(1) if h.strip()]

HEADERS = [
    # your form fields + demos...
    "Business Name","Location Address","Location Type","Years in Operation","Days/Hours of Operation",
    "Seating Indoor","Seating Outdoor","Full Kitchen","Average Weekly Sales","Average Transaction Value",
    "Top Item 1","Top Item 2","Top Item 3","% Sales Drinks","% Sales Food","Peak Hours",
    "Sales Monday","Sales Tuesday","Sales Wednesday","Sales Thursday","Sales Friday","Sales Saturday","Sales Sunday",
    "Avg Transactions/Day","Mobile/Online Ordering","Avg Wait Time","Target Age Range","Main Customer Segments",
    "% Regulars","Collect Contact Info","Run Surveys","# Competitors within 1mi","Competitor 1","Competitor 2","Competitor 3",
    "Near POI","Marketing Strategy","Partner Local","# Employees","Full-Time Staff","Part-Time Staff",
    "Track Labor % of Revenue","Staffing Challenges","Training Procedures","POS Export Filename","Labor Report Filename","Benchmark Similar Businesses",
    "Additional Insight","Timestamp",
    # demographic keys...
    "Total Population (2024)", "Total Households (2024)", "Avg Household Size (2024)", "Median Household Income (2024)",
    "Per-Capita Income (2024)", "Diversity Index (2024)", "Pop. Growth Rate ’24–’29 (%)", "Median Home Value (2024)",
    "Avg Home Value (2024)", "Daytime Worker Pop. (2024)", "Median Age (2024)", "Owner-occupied Households (2024)",
    "Renter-occupied Households (2024)", "Vacant Housing Units (2024)", "Daytime Population (2024)", "Daytime Resident Pop. (2024)",
    "Population Age 18 (2024)", "Population Age 25 (2024)", "Population Age 35 (2024)", "Population Age 65+ (2024)",
    "Age Dependency Ratio (2024)", "White Population (2024)", "Black Population (2024)", "Asian Population (2024)",
    "Unemployment Rate (2024)", "Household Growth Rate ’24–’29 (%)", "Eating & Drinking Businesses (2024)", "Group Quarters Population (2024)"
]

_existing = get_headers()
if _existing != HEADERS:
    # wipe & rewrite header row once
    ws.delete_rows(1)
    ws.insert_row(HEADERS, 1)

def append_row_to_sheet(row: list):
    ws.append_row(row, value_input_option="USER_ENTERED")

# Cache the entire sheet as a DataFrame (refresh every 5 minutes)

def load_sheet_df() -> pd.DataFrame:
    values = ws.get_all_values()
    if not values:
        return pd.DataFrame()
    raw_headers = values[0]
    idxs = [i for i,h in enumerate(raw_headers) if h.strip()]
    headers = [raw_headers[i] for i in idxs]
    data = []
    for row in values[1:]:
        row += [""] * (len(raw_headers) - len(row))
        data.append([row[i] for i in idxs])
    return pd.DataFrame(data, columns=headers)

# Save demographics helper
def save_demographics_to_sheet(business_name: str, demo: dict):

    df = load_sheet_df()
    try:
        # locate the row by Business Name
        row_idx = df.index[df["Business Name"] == business_name][0] + 2
    except IndexError:
        st.error(f"Business '{business_name}' not found in sheet!")
        return

    # Fetch current headers
    headers = [h for h in ws.row_values(1) if h.strip() != ""]

    # Add any new demo keys as columns
    new_cols = [k for k in demo if k not in headers]
    if new_cols:
        headers += new_cols
        ws.update("A1:1", [headers])

    # Write each demo value into its cell
    with st.spinner("Saving demographics…"):
        prog = st.progress(0)
        total = len(demo)
        for i, (k, v) in enumerate(demo.items(), start=1):
            col_idx = headers.index(k) + 1
            ws.update_cell(row_idx, col_idx, str(v))
            prog.progress(i / total)

    st.success("Demographics saved!")


# --- STREAMLIT APP ---
def main():
    st.title("Take My Advice [...]")
    tab_form, tab_results, tab_explore, tab_optimize, tab_deep ,tab_summary = st.tabs(
        ["Form","Results","Exploration","Optimizer", "Deep Dive", "Summary"]
    )

    # --- Form Tab ---
    with tab_form:
        st.header("CAFÉ PROFILE")
        business_name = st.text_input("Business Name")
        location_address = st.text_input("Location Address")
        location_type = st.radio("Standalone or inside another business?", ["Standalone","Inside another business"])
        years_in_operation = st.number_input("Years in operation", 0, step=1)
        days_hours = st.text_area("Days/hours of operation", help="e.g. Mon–Fri 7:00–19:00; Sat 8:00–16:00")
        seating_indoor = st.number_input("Seating capacity (indoor)", 0, step=1)
        seating_outdoor = st.number_input("Seating capacity (outdoor)", 0, step=1)
        full_kitchen = st.radio("Do you have a full kitchen?", ["Yes","No"])

        st.header("SALES & POS DATA")
        avg_weekly_sales = st.number_input("Average weekly sales (last 4 weeks)", 0.0, format="%.2f")
        avg_transaction_value = st.number_input("Average transaction value", 0.0, format="%.2f")
        st.subheader("Top 3 best-selling items")
        top1 = st.text_input("1st best-selling item")
        top2 = st.text_input("2nd best-selling item")
        top3 = st.text_input("3rd best-selling item")
        pct_drinks = st.slider("% of sales from drinks", 0,100,50)
        pct_food = 100 - pct_drinks
        peak_hours = st.multiselect("Peak hours (POS)", [f"{h}:00" for h in range(24)])
        st.subheader("Sales by day of week")
        days = ["Monday","Tuesday","Wednesday","Thursday","Friday","Saturday","Sunday"]
        sales_by_day = {d: st.number_input(d, 0.0, format="%.2f") for d in days}
        avg_transactions = st.number_input("# transactions/day (avg)", 0, step=1)
        online_sales = st.radio("Mobile/online ordering?", ["Yes","No"])
        avg_wait_time = st.number_input("Avg wait time (min)", 0.0, format="%.1f")

        st.header("CUSTOMER DEMOGRAPHICS & INSIGHT")
        age_range = st.text_input("Target customer age range")
        segments = st.multiselect("Main customer segments", ["Students","Professionals","Tourists","Others"])
        pct_regulars = st.slider("% regular customers (3+ visits/mo)", 0,100,20)
        collect_emails = st.radio("Collect customer contacts?", ["Yes","No"])
        run_surveys = st.radio("Run surveys/feedback?", ["Yes","No"])

        st.header("COMPETITION & LOCATION")
        num_competitors = st.number_input("# competitors within 1 mile", 0, step=1)
        comp1 = st.text_input("Competitor 1 (name+loc)")
        comp2 = st.text_input("Competitor 2 (name+loc)")
        comp3 = st.text_input("Competitor 3 (name+loc)")
        near_poi = st.radio("Near major POI?", ["Yes","No"])
        marketing = st.selectbox("Marketing strategy", ["Social media","Events","None","Other"])
        partner_local = st.radio("Partner with local businesses/events?", ["Yes","No"])

        st.header("OPERATIONS & STAFFING")
        num_employees = st.number_input("Number of employees", 0, step=1)
        full_time = st.number_input("# full-time staff", 0, step=1)
        part_time = st.number_input("# part-time staff", 0, step=1)
        track_labor_pct = st.radio("Track labor % of revenue?", ["Yes","No"])
        staffing_challenges = st.text_area("Staffing challenges")
        training_procedures = st.radio("Standard training procedures?", ["Yes","No"])


        st.header("OPTIONAL TECHNICAL INPUTS")
        pos_export   = st.file_uploader("POS export (CSV)", type=["csv"])

        if pos_export is not None:
           st.markdown("✅ POS file loaded.")
           if st.button("Upload POS to sheet"):
               # Read the CSV
               df_pos = pd.read_csv(pos_export)
               # Try to open or else create the "POS Raw" sheet
               try:
                   ws_pos = gc.open_by_key(SHEET_ID).worksheet("POS Raw")
               except gspread.WorksheetNotFound:
                   ws_pos = gc.open_by_key(SHEET_ID).add_worksheet(
                       title="POS Raw",
                       rows=str(len(df_pos) + 10),
                       cols=str(len(df_pos.columns) + 5),
                   )
            # Clear old data and write new
               ws_pos.clear()
               set_with_dataframe(ws_pos, df_pos)
               st.success("POS CSV uploaded to tab ‘POS Raw’")
        labor_report = st.file_uploader("Employee schedule/report", type=["csv","xlsx","pdf"])
        benchmark    = st.radio("Benchmark against similar businesses?", ["Yes","No"])

                # NEW: freeform Additional Insight
        additional_insight = st.text_area(
            "Additional insight (any extra context for the Deep Dive)",
            help="This will be saved alongside your form inputs and fed into the Deep Dive model."
        )

        if st.button("Submit"):
            row = [
                business_name, location_address, location_type, years_in_operation, days_hours,
                seating_indoor, seating_outdoor, full_kitchen,
                avg_weekly_sales, avg_transaction_value,
                top1, top2, top3, pct_drinks, pct_food, ",".join(peak_hours),
                sales_by_day["Monday"], sales_by_day["Tuesday"], sales_by_day["Wednesday"],
                sales_by_day["Thursday"], sales_by_day["Friday"], sales_by_day["Saturday"],
                sales_by_day["Sunday"], avg_transactions, online_sales, avg_wait_time,
                age_range, ",".join(segments), pct_regulars, collect_emails, run_surveys,
                num_competitors, comp1, comp2, comp3, near_poi, marketing, partner_local,
                num_employees, full_time, part_time, track_labor_pct, staffing_challenges, training_procedures,
                pos_export.name if pos_export else "", labor_report.name if labor_report else "", benchmark,
                additional_insight,                     # ← added here
                datetime.utcnow().isoformat()
            ]
            append_row_to_sheet(row)
            st.success("Form submitted!")

    # --- Results Tab ---
    with tab_results:
        st.header("Foot-Traffic Model Results")
        res = run_foot_traffic()
        if res["mae"] is not None:
            st.metric("MAE", f"{res['mae']:.2f}")
            st.pyplot(res["figure"])
        else:
            st.info("Not enough data yet.")

        st.header("Margin Optimization Results")
        mres = run_margin()
        if mres["metrics"]:
            st.metric("MAE", f"{mres['metrics']['MAE']:.2f}")
            st.metric("R²", f"{mres['metrics']['R2']:.2f}")
            st.pyplot(mres["figure"])
        else:
            st.info("Not enough data yet.")

    # --- Exploration Tab ---
    with tab_explore:
        st.header("Exploration")

        df = load_sheet_df()
        if df.empty:
            st.info("No cafés on file yet.")
        else:
            cafelist = df["Business Name"].unique().tolist()
            selected = st.selectbox("Select café to explore", cafelist)

            record = df[df["Business Name"] == selected].iloc[0]
            address = record["Location Address"]

            st.markdown(f"**Address:** {address}")

            if st.button("Explore"):
                map_html, demo, cafes = run_explorer(address)
                if map_html:
                    st.session_state["last_demo"] = demo
                    st.session_state["last_cafes"] = cafes
                    st.session_state["target_cafe_name"] = selected
                    st.session_state["target_cafe_address"] = address

                    st.subheader("Map View")
                    st.components.v1.html(map_html, height=450)

                    st.subheader("Key Demographics")
                    df_demo = pd.DataFrame(list(demo.items()), columns=["Description", "Value"])
                    st.table(df_demo)

                    st.subheader("Nearby Cafés")
                    df_cafes = pd.DataFrame(cafes)
                    for c in df_cafes.columns:
                        df_cafes[c] = df_cafes[c].fillna("").astype(str)
                    st.session_state["last_cafes_df"] = df_cafes
                    st.table(df_cafes)
                else:
                    st.error("Map generation failed.")

            if "last_demo" in st.session_state and st.button("Save demographics"):
                with st.spinner("Saving demographics to sheet…"):
                    save_demographics_to_sheet(
                        selected,
                        st.session_state["last_demo"]
                    )

            # Compare with similar regions
            if "last_demo" in st.session_state:
                st.subheader("Compare with Similar Demographic Regions")
                city = st.text_input("Enter city to search for similar regions", "Los Angeles")
                k = st.slider("How many similar regions to find?", 1, 10, 5)

                if st.button("Find Similar Regions"):
                    token = get_token()
                    lat, lon = geocode(city, token)
                    if lat is not None:
                        batch = fetch_demographics_batch(lat, lon, token, n=15)
                        top_k = get_similar_regions(st.session_state["last_demo"], batch, k)
                        st.session_state["top_k_regions"] = top_k

                        region_names = top_k["_region_name"].tolist()
                        selected_region = st.selectbox("Choose a similar region", region_names)

                        selected_row = top_k[top_k["_region_name"] == selected_region].iloc[0]
                        region_lat, region_lon = selected_row["_lat"], selected_row["_lon"]

                        cafes_nearby = fetch_cafes_in_region(region_lat, region_lon)
                        st.session_state["similar_region_cafes"] = cafes_nearby

                        st.subheader("Cafés in Selected Similar Region")
                        df_region_cafes = pd.DataFrame(cafes_nearby)
                        st.session_state["region_cafe_df"] = df_region_cafes
                        st.table(df_region_cafes)
                    else:
                        st.warning("Could not geocode city.")

            if "region_cafe_df" in st.session_state:
                        st.subheader("Compare with a Café from Selected Region")
                        options = st.session_state["region_cafe_df"]["Name"].tolist()
                        selected_comp = st.selectbox("Choose a competitor café", options)

                        comp_row = st.session_state["region_cafe_df"][st.session_state["region_cafe_df"]["Name"] == selected_comp].iloc[0]

                        # Construct GPT prompt
                        prompt = f"""
                        The target café is {st.session_state['target_cafe_name']} located at {st.session_state['target_cafe_address']}.
                        Here are the key demographics: {st.session_state['last_demo']}.

                        The competitor café is {comp_row['Name']} located at {comp_row['Address']}.
                        It appears in a nearby region with similar demographics.

                        Please compare these two cafés from a business strategy perspective:
                        - Who they likely attract?
                        - Opportunities for the target café to compete or differentiate?
                        - What the competitor might be doing well?
                        Return this in bullet points. Ensure each point has examples pulled from the internet or data to provide meaninfgul comparison
                        """

                        response = client.chat.completions.create(
                            model="gpt-4o",
                            messages=[{"role": "user", "content": prompt}]
                        )

                        st.subheader("GPT Insights")
                        st.markdown(response.choices[0].message.content)

    # --- Optimizer Tab ---
    with tab_optimize:
        st.header("🔧 Optimization Assistant")

        # 1) Load submissions
        df = load_sheet_df()
        if df.empty:
            st.info("No data submitted yet.")
            st.stop()

        # 2) Café selector
        cafes = df["Business Name"].unique().tolist()
        selected = st.selectbox("Select café", cafes, key="optimize_cafe")
        record = df[df["Business Name"] == selected].iloc[0]

        # 3) Split form vs. demographics
        headers = get_headers()
        ts_index = headers.index("Timestamp") + 1
        form_data = {c: record[c] for c in headers[:ts_index]}
        demo_data = {c: record[c] for c in headers[ts_index:] if pd.notna(record[c])}

        # 4) Load POS
        ws_pos = get_pos_sheet()
        pos_data = (
            pd.DataFrame(ws_pos.get_all_records()).to_dict("records")
            if ws_pos else []
        )

        extra = record.get("Additional Insight", "")
        context = {
            "form": form_data,
            "demographics": demo_data,
            "pos": pos_data,
            "pos_raw": pos_data,
            "labor_raw": None,
            "extra": extra,
        }

        if st.button("Run Optimization Engine"):
            with st.spinner("Running GPT + ML to optimize your café operations..."):
                results = run_optimization(context)

            st.success("Optimization complete!")
            for res in results["responses"]:
                st.subheader(f"🧠 Prompt: {res['prompt']}")
                st.markdown(res["answer"])

            if "charts" in results:
                st.markdown("---")
                st.subheader("📊 Visual Insights")

                if "prophet" in results["charts"]:
                    st.markdown("**Sales Forecast (Prophet)**")
                    st.pyplot(results["charts"]["prophet"])

                if "kmeans" in results["charts"]:
                    st.markdown("**Menu Clustering (KMeans)**")
                    st.pyplot(results["charts"]["kmeans"])

                if "regression" in results["charts"]:
                    st.markdown("**Labor Efficiency (RandomForest)**")
                    st.pyplot(results["charts"]["regression"])

                if "xgboost" in results["charts"]:
                    st.markdown("**Waste Drivers (XGBoost Feature Importance)**")
                    st.pyplot(results["charts"]["xgboost"])

        st.markdown("---")
        st.subheader("Ask the Optimizer Directly")
        custom_q = st.text_area("Enter a specific optimization question")
        if st.button("Ask Optimizer"):
            with st.spinner("Analyzing your custom request..."):
                follow_up = run_optimization(context, prompts=[custom_q])
            for res in follow_up["responses"]:
                st.subheader(res["prompt"])
                st.write(res["answer"])

    # --- Presentation Tab ---
    with tab_summary:
        st.header("Interactive Presentation")
        st.write("🧪 Presentation tab loaded correctly.")

        # Example demo content (replace with real summaries later)
        sections = [
            ("Cafe Summary", ["Eruta Nature is a nature-themed café.", "Located in Los Angeles."]),
            ("Key Demographics", ["Young adult population", "High diversity index"]),
            ("Explore Insights", ["Strong regular base", "Bagel combo performs well"]),
            ("Optimizer Suggestions", ["Reduce spend on underperforming items", "Optimize labor scheduling"]),
            ("Deep Dive Goals", ["Align product mix with demographics", "Test targeted offers"]),
            ("Next Steps", ["Pilot new menu items", "Monitor week-on-week revenue change"])
        ]

        if st.button("Generate Presentation"):
            progress = st.progress(0.0)
            progress.progress(0.3)

            html = generate_impress_html(sections)
            progress.progress(0.8)

            components.html(html, height=600, scrolling=False)
            progress.progress(1.0)


    # --- Deep Dive Tab ---
    with tab_deep:
        st.header("Deep Dive Insights")

        # 1) Load submissions
        df = load_sheet_df()
        if df.empty:
            st.info("No data submitted yet.")
            st.stop()

        # 2) Café selector
        cafes = df["Business Name"].unique().tolist()
        selected = st.selectbox("Select café", cafes)
        record = df[df["Business Name"] == selected].iloc[0]

        # 3) Split form vs. demographics
        headers = get_headers()
        ts_index = headers.index("Timestamp") + 1
        form_data = {c: record[c] for c in headers[:ts_index]}
        demo_data = {c: record[c] for c in headers[ts_index:] if pd.notna(record[c])}

        # 4) Load POS if exists
        ws_pos = get_pos_sheet()
        pos_data = (
            pd.DataFrame(ws_pos.get_all_records()).to_dict("records")
            if ws_pos
            else []
        )

        extra = record.get("Additional Insight", "")

        # 5) Pick which goals
        st.markdown("### Pick which Deep-Dive areas to run:")
        do_marketing = st.checkbox("🟦 Marketing", value=True)
        do_finance   = st.checkbox("🟥 Finance", value=True)
        do_compete   = st.checkbox("🟩 Competitor Insight", value=True)

        # 6) Define question buckets
        marketing_qs = [
            "Based on local ArcGIS demographics (age, income, ethnicity), do your current top-selling items align with dominant consumer preferences?",
            "What percentage of your highest AOV customers belong to the reported target group (e.g., 18–24, students)?",
            "Are there underserved demographic groups nearby (e.g., young professionals, working parents, older adults) who might respond to different offerings?",
        ]
        finance_qs = [
            "Which top 10 selling items yield the highest vs. lowest net profit margins, after accounting for modifiers, discounts, and waste?",
            "Are there high-volume, low-margin items (e.g., bagels, basic drinks) that could be bundled or upsold to increase AOV?",
            "What % of weekly revenue comes from your top 3 items (e.g., Bagel Combo, Latte, Matcha)?",
            "How do average order values differ between weekdays and weekends, and are certain hours consistently underperforming?",
            "Based on item-level waste and refund rates, which products contribute most to unnecessary cost?",
            "How many orders are fulfilled per employee during each shift, and where are the biggest bottlenecks?",
        ]
        competitor_qs = [
            "How does your demographic profile (students 18–24) compare to those of your top 3 competitors within 1 mile?",
            "Which competitors are leveraging seasonal specials, mobile ordering, or loyalty programs to drive repeat traffic — and what could you adapt?",
            "Which product types (e.g., specialty drinks, protein-rich meals) are you not offering that competitors succeed with?",
            "Is your revenue/sqft higher or lower than cafés with similar sales volume, seating, and staff count?",
            "Which of your current differentiators (e.g., indoor space, menu creativity, healthiness) are under-leveraged based on online comparisons?",
        ]

        queries = []
        if do_marketing:
            queries += marketing_qs
        if do_finance:
            queries += finance_qs
        if do_compete:
            queries += competitor_qs

        if not queries:
            st.warning("Select at least one area to deep-dive on.")
        else:
            context = {
                "form": form_data,
                "demographics": demo_data,
                "pos": pos_data,
                "extra": extra,
            }
            deep = run_deep_dive(context, queries)

            for resp in deep["responses"]:
                st.subheader(resp["query"])
                st.write(resp["answer"])

            if st.button("📑 Generate PPT & PDF Report"):
                from pptx import Presentation
                from pptx.util import Pt, Inches
                from pptx.dml.color import RGBColor
                from pptx.enum.shapes import MSO_SHAPE
                from reportlab.pdfgen import canvas
                from reportlab.lib.pagesizes import letter
                import io
                from datetime import datetime
                import math

                prog = st.progress(0)
                grouped = {}
                for resp in deep["responses"]:
                    q, a = resp["query"], resp["answer"]
                    if q in marketing_qs:
                        goal = "Marketing"
                    elif q in finance_qs:
                        goal = "Finance"
                    else:
                        goal = "Competitor Insight"
                    grouped.setdefault(goal, []).append((q, a))
                prog.progress(0.3)

                prs = Presentation()
                theme_gold = RGBColor(189, 169, 103)
                theme_text = RGBColor(176, 166, 152)
                theme_bg = RGBColor(18, 18, 18)

                title_slide = prs.slides.add_slide(prs.slide_layouts[6])
                shape = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
                shape.fill.solid()
                shape.fill.fore_color.rgb = theme_bg

                title = title_slide.shapes.title or title_slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
                title_tf = title.text_frame
                title_tf.clear()
                p = title_tf.paragraphs[0]
                p.text = f"Take my Advice – {selected}"
                p.font.size = Pt(36)
                p.font.bold = True
                p.font.color.rgb = theme_gold

                sub = title_slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1)).text_frame
                p = sub.paragraphs[0]
                p.text = "Generated on " + datetime.utcnow().strftime("%Y-%m-%d")
                p.font.size = Pt(18)
                p.font.color.rgb = theme_text

                all_qas = [(goal, q, a) for goal, qas in grouped.items() for q, a in qas]
                total_qas = len(all_qas)
                slides_needed = 10
                chunk_size = math.ceil(total_qas / slides_needed)
                chunks = [all_qas[i:i + chunk_size] for i in range(0, total_qas, chunk_size)]
                while len(chunks) < slides_needed:
                    chunks.append([])

                for chunk in chunks:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = theme_bg

                    tf = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6)).text_frame
                    tf.word_wrap = True
                    for goal, q, a in chunk:
                        p1 = tf.add_paragraph()
                        p1.text = f"[{goal}]"
                        p1.font.size = Pt(16)
                        p1.font.bold = True
                        p1.font.color.rgb = theme_gold

                        p2 = tf.add_paragraph()
                        p2.text = f"Q: {q}"
                        p2.font.size = Pt(14)
                        p2.font.bold = True
                        p2.font.color.rgb = theme_text

                        p3 = tf.add_paragraph()
                        p3.text = f"A: {a}"
                        p3.font.size = Pt(14)
                        p3.font.color.rgb = theme_text

                        tf.add_paragraph().text = ""

                pptx_io = io.BytesIO()
                prs.save(pptx_io)
                pptx_bytes = pptx_io.getvalue()
                prog.progress(0.6)

                pdf_io = io.BytesIO()
                c = canvas.Canvas(pdf_io, pagesize=letter)
                width, height = letter
                y = height - 50
                c.setFont("Helvetica-Bold", 18)
                c.setFillColorRGB(189/255, 169/255, 103/255)
                c.drawString(50, y, f"Café Deep Dive – {selected}")
                y -= 40
                for goal, items in grouped.items():
                    c.setFont("Helvetica-Bold", 14)
                    c.setFillColorRGB(1, 1, 1)
                    c.drawString(50, y, f"{goal} Insights")
                    y -= 20
                    c.setFont("Helvetica", 12)
                    for q, a in items:
                        lines = [f"Q: {q}", f"A: {a}"]
                        for line in lines:
                            for subline in line.split("\n"):
                                c.drawString(70, y, subline)
                                y -= 15
                                if y < 50:
                                    c.showPage()
                                    y = height - 50
                                    c.setFont("Helvetica", 12)
                    y -= 20
                c.save()
                pdf_bytes = pdf_io.getvalue()
                prog.progress(1.0)

                st.download_button("Download PPTX", pptx_bytes, file_name=f"Take My Advice [{selected}].pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
                st.download_button("Download PDF", pdf_bytes, file_name=f"Take My Advice [{selected}].pdf", mime="application/pdf")

        st.markdown("---")
        st.subheader("Chat with the Café AI")
        user_q = st.text_area("Your question")
        if st.button("Send"):
            with st.spinner("🤖 Café AI is thinking…"):
                follow = run_deep_dive(context, [user_q])
            for resp in follow["responses"]:
                st.subheader(resp["query"])
                st.write(resp["answer"])


if __name__ == "__main__":
    main()
